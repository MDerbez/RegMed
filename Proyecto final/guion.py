from pptx import Presentation

# Cargar la presentación
prs = Presentation("pantallas.pptx")

# Crear archivo de salida
with open("notas.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides, start=1):
        # Verificar si la diapositiva tiene notas
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        else:
            notes = ""
        
        f.write(f"Diapositiva {i}:\n")
        f.write(notes + "\n\n")

print("✅ Notas extraídas en 'notas.txt'")
